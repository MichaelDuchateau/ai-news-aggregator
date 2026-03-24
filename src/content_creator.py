"""Content creator for generating presentation slides."""

from typing import List
from pathlib import Path
import subprocess
import json

from src.models import NewsItem
from src.config import Config


class ContentCreator:
    """Creates presentation slides for deep dive items."""
    
    def __init__(self, config: Config):
        self.config = config
        self.presentation_style = config.get('presentation.style', 'minimal')
    
    def create_presentation(self, items: List[NewsItem], week: str, output_dir: Path) -> str:
        """
        Create a presentation from deep dive items.
        Returns the path to the created file.
        """
        print(f"\n🎯 Creating presentation for {len(items)} items...")
        
        # Prepare output
        output_dir.mkdir(parents=True, exist_ok=True)
        output_file = output_dir / f"AI-News-{week}.pptx"
        
        # Create HTML slides
        print("  📝 Generating slides...")
        slides_dir = output_dir / "slides_html"
        slides_dir.mkdir(exist_ok=True)
        
        self._create_title_slide(slides_dir, week)
        
        for i, item in enumerate(items, 1):
            self._create_item_slide(slides_dir, item, i)
            item.slide_file = output_file.name
            item.slide_number = i
        
        # Convert to PPTX using html2pptx
        print("  🔄 Converting to PowerPoint...")
        success = self._convert_to_pptx(slides_dir, output_file)
        
        if success:
            print(f"  ✅ Presentation created: {output_file.name}")
            return str(output_file)
        else:
            print("  ⚠️  Error creating presentation")
            return None
    
    def _create_title_slide(self, slides_dir: Path, week: str):
        """Create title slide."""
        title_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <link rel="stylesheet" href="styles.css">
</head>
<body>
    <div class="slide title-slide">
        <div class="title-content">
            <h1>AI News</h1>
            <h2>Week {week}</h2>
            <p class="subtitle">Key Developments in Artificial Intelligence</p>
        </div>
    </div>
</body>
</html>"""
        
        with open(slides_dir / "slide_0_title.html", 'w') as f:
            f.write(title_html)
    
    def _create_item_slide(self, slides_dir: Path, item: NewsItem, number: int):
        """Create a minimal, narration-ready slide for an item."""
        
        # Extract key visual element or quote from narrative
        visual_text = self._extract_key_point(item)
        
        # Minimal slide following Garr Reynolds philosophy
        slide_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <link rel="stylesheet" href="styles.css">
</head>
<body>
    <div class="slide content-slide">
        <div class="slide-content">
            <h1>{item.title}</h1>
            <div class="visual-element">
                {visual_text}
            </div>
            <div class="source">{item.source}</div>
        </div>
    </div>
    <div class="notes">
        <h3>Narrative for Narration:</h3>
        <p>{item.narrative}</p>
        
        <h3>Source:</h3>
        <p><a href="{item.url}">{item.url}</a></p>
        
        <h3>Tags:</h3>
        <p>{', '.join(item.tags)}</p>
    </div>
</body>
</html>"""
        
        with open(slides_dir / f"slide_{number}.html", 'w') as f:
            f.write(slide_html)
    
    def _extract_key_point(self, item: NewsItem) -> str:
        """Extract a key point or quote for the slide visual."""
        if not item.narrative:
            return f"<p class='key-point'>{item.summary[:150]}...</p>"
        
        # Extract first compelling sentence or key finding
        sentences = item.narrative.split('. ')
        for sentence in sentences[:3]:
            # Look for sentences with impact keywords
            if any(word in sentence.lower() for word in 
                   ['breakthrough', 'first', 'new', 'announced', 'released', 'achieved']):
                return f"<p class='key-point'>{sentence}.</p>"
        
        # Default to first substantial sentence
        if len(sentences) > 0:
            return f"<p class='key-point'>{sentences[0]}.</p>"
        
        return f"<p class='key-point'>{item.title}</p>"
    
    def _create_styles(self, slides_dir: Path):
        """Create CSS for minimal, visual slides."""
        css = """/* Minimal AI News Slides - Garr Reynolds Philosophy */

:root {
    --primary-color: #667eea;
    --secondary-color: #764ba2;
    --text-dark: #2d3748;
    --text-light: #718096;
    --background: #ffffff;
}

* {
    margin: 0;
    padding: 0;
    box-sizing: border-box;
}

body {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
    background: var(--background);
}

.slide {
    width: 960px;
    height: 540px;
    display: flex;
    align-items: center;
    justify-content: center;
    padding: 60px;
}

/* Title Slide */
.title-slide {
    background: linear-gradient(135deg, var(--primary-color) 0%, var(--secondary-color) 100%);
    color: white;
    text-align: center;
}

.title-content h1 {
    font-size: 72px;
    font-weight: 700;
    margin-bottom: 20px;
    letter-spacing: -2px;
}

.title-content h2 {
    font-size: 48px;
    font-weight: 300;
    margin-bottom: 30px;
}

.subtitle {
    font-size: 24px;
    opacity: 0.9;
    font-weight: 300;
}

/* Content Slide */
.content-slide {
    flex-direction: column;
    justify-content: space-between;
    align-items: flex-start;
}

.slide-content {
    width: 100%;
}

.slide-content h1 {
    font-size: 48px;
    color: var(--text-dark);
    margin-bottom: 60px;
    line-height: 1.2;
    font-weight: 600;
}

.visual-element {
    margin: 60px 0;
}

.key-point {
    font-size: 36px;
    color: var(--primary-color);
    line-height: 1.4;
    font-weight: 400;
}

.source {
    font-size: 20px;
    color: var(--text-light);
    margin-top: auto;
    padding-top: 40px;
    border-top: 2px solid #e2e8f0;
}

/* Speaker Notes */
.notes {
    display: none;
}
"""
        
        with open(slides_dir / "styles.css", 'w') as f:
            f.write(css)
    
    def _convert_to_pptx(self, slides_dir: Path, output_file: Path) -> bool:
        """Convert HTML slides to PPTX using python-pptx."""
        try:
            # Create CSS file
            self._create_styles(slides_dir)
            
            # Use python-pptx for a simpler approach
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            
            # Get HTML files in order
            html_files = sorted(slides_dir.glob("slide_*.html"))
            
            for html_file in html_files:
                # Parse HTML to extract content
                content = self._parse_slide_html(html_file)
                
                if 'title-slide' in str(html_file):
                    self._add_title_slide_pptx(prs, content)
                else:
                    self._add_content_slide_pptx(prs, content)
            
            # Save
            prs.save(str(output_file))
            return True
            
        except Exception as e:
            print(f"    Error converting to PPTX: {e}")
            return False
    
    def _parse_slide_html(self, html_file: Path) -> dict:
        """Parse HTML slide to extract content."""
        from bs4 import BeautifulSoup
        
        with open(html_file, 'r') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        
        content = {}
        
        h1 = soup.find('h1')
        if h1:
            content['title'] = h1.get_text()
        
        h2 = soup.find('h2')
        if h2:
            content['subtitle'] = h2.get_text()
        
        key_point = soup.find('p', class_='key-point')
        if key_point:
            content['key_point'] = key_point.get_text()
        
        subtitle_p = soup.find('p', class_='subtitle')
        if subtitle_p:
            content['tagline'] = subtitle_p.get_text()
        
        source = soup.find('div', class_='source')
        if source:
            content['source'] = source.get_text()
        
        notes = soup.find('div', class_='notes')
        if notes:
            content['notes'] = notes.get_text(separator='\n')
        
        return content
    
    def _add_title_slide_pptx(self, prs, content):
        """Add title slide to presentation."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add gradient background (simplified - solid color)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(102, 126, 234)
        
        # Title
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'AI News')
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if 'subtitle' in content:
            sub_top = Inches(3)
            sub_box = slide.shapes.add_textbox(left, sub_top, width, Inches(1))
            sub_frame = sub_box.text_frame
            sub_frame.text = content['subtitle']
            
            p = sub_frame.paragraphs[0]
            p.font.size = Pt(48)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        
        # Tagline
        if 'tagline' in content:
            tag_top = Inches(4.2)
            tag_box = slide.shapes.add_textbox(left, tag_top, width, Inches(0.5))
            tag_frame = tag_box.text_frame
            tag_frame.text = content['tagline']
            
            p = tag_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide_pptx(self, prs, content):
        """Add content slide to presentation."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1.2)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', '')
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(45, 55, 72)
        
        # Key point
        if 'key_point' in content:
            kp_top = Inches(2.2)
            kp_height = Inches(2)
            kp_box = slide.shapes.add_textbox(left, kp_top, width, kp_height)
            kp_frame = kp_box.text_frame
            kp_frame.text = content['key_point']
            kp_frame.word_wrap = True
            
            p = kp_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.color.rgb = RGBColor(102, 126, 234)
            p.line_spacing = 1.3
        
        # Source
        if 'source' in content:
            src_top = Inches(4.8)
            src_box = slide.shapes.add_textbox(left, src_top, width, Inches(0.5))
            src_frame = src_box.text_frame
            src_frame.text = content['source']
            
            p = src_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(113, 128, 150)
        
        # Speaker notes
        if 'notes' in content:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            notes_frame.text = content['notes']
